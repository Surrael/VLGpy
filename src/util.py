from typing import Literal, List
from openai import OpenAI
import os
import fitz
import pyttsx3
from pptx import Presentation



import ffmpeg

current_dir = os.getcwd()
base_dir = os.path.dirname(current_dir)
os.chdir(base_dir)


def pdf_to_images(pdf_path: str, output_folder: str = "dir") -> list:
    """
    Convert each page of a PDF into an image and save them in the specified output folder.

    Args:
        pdf_path (str): Path to the PDF file.
        output_folder (str): Output folder to save the images. Defaults to "dir".

    Returns:
        list: List of paths to the generated images.
    """
    image_paths = []

    # Open the PDF file
    pdf_document = fitz.open(pdf_path)

    # Iterate through each page in the PDF
    for page_number in range(len(pdf_document)):
        # Get the page
        page = pdf_document.load_page(page_number)

        # Render the page to an image
        image = page.get_pixmap()

        # Save the image
        image_path = f"{output_folder}/page_{page_number + 1}.png"
        image.save(image_path)

        print(f"Page {page_number} saved as {image_path}")

        # Append the path to the list of image paths
        image_paths.append(image_path)

    # Close the PDF document
    pdf_document.close()

    return image_paths


def extract_pptx_notes(path: str) -> list:
    ppt = Presentation(path)

    notes = []

    for page, slide in enumerate(ppt.slides):
        text = slide.notes_slide.notes_text_frame.text
        notes.append(text)

    return notes


def text_to_speech(texts: List[str], voice: Literal["alloy", "echo", "fable", "onyx", "nova", "shimmer"], key: str,
                   path: str = "dir") -> List[str]:
    """
    Convert a list of texts to speech using the specified TTS voice and save each audio as an MP3 file.

    Args:
        texts (List[str]): The list of texts to be converted to speech.
        voice (Literal["alloy", "echo", "fable", "onyx", "nova", "shimmer"]): The TTS voice to be used.
        path (str, optional): The directory where the audio files will be saved. Defaults to "dir".

    Returns:
        List[str]: List of paths to the generated MP3 files.
    """
    client = OpenAI(api_key=key)

    mp3_paths = []

    for i, text in enumerate(texts):
        response = client.audio.speech.create(
            model="tts-1",
            voice=voice,
            input=text
        )
        mp3_path = f"{path}/audio_{i}.mp3"
        response.write_to_file(mp3_path)
        mp3_paths.append(mp3_path)

    return mp3_paths


def text_to_speech_demo(texts: List[str], path: str = "dir") -> List[str]:
    engine = pyttsx3.init()
    mp3_paths = []

    for i, text in enumerate(texts):
        mp3_path = f"{path}/audio_{i}.mp3"
        engine.save_to_file(text, mp3_path)
        engine.runAndWait()
        mp3_paths.append(mp3_path)

    engine.stop()
    return mp3_paths


def parse_script_file(script_path: str) -> list:
    """
    Parse a .txt script file for video slides.

    Args:
        script_path (str): Path to the script file.

    Returns:
        list: List of slide texts.
    """
    slides = []
    with open(script_path, 'r') as file:
        current_slide_text = ''
        for line in file:
            line = line.strip()
            if line == '#NEXT':
                slides.append(current_slide_text.strip())
                current_slide_text = ''
            else:
                current_slide_text += line + ' '
        # Append the last slide text
        slides.append(current_slide_text.strip())
    return slides
